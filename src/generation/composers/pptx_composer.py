"""PowerPoint (PPTX) composer for generating presentations from pitch content."""

from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .base import BaseComposer, ComposerConfig, ComposerResult, ThemeColor
from ...models.pitch import Pitch, PitchSection, SectionType


@dataclass
class PPTXConfig(ComposerConfig):
    """Configuration for PPTX composer."""

    # Slide dimensions (default: widescreen 16:9)
    slide_width: float = 13.333  # inches
    slide_height: float = 7.5  # inches

    # Font settings
    title_font_name: str = "Calibri"
    title_font_size: int = 44
    subtitle_font_size: int = 24
    body_font_name: str = "Calibri"
    body_font_size: int = 18
    bullet_font_size: int = 16

    # Layout settings
    margin_left: float = 0.5  # inches
    margin_right: float = 0.5
    margin_top: float = 0.75
    margin_bottom: float = 0.5

    # Content settings
    max_bullets_per_slide: int = 6
    include_talking_points_slide: bool = False

    # Slide options
    add_slide_numbers: bool = True
    add_title_footer: bool = True


class PPTXComposer(BaseComposer):
    """Composer for generating PowerPoint presentations from pitch content."""

    def __init__(self, config: Optional[PPTXConfig] = None):
        super().__init__(config or PPTXConfig())
        self.config: PPTXConfig = self.config  # type: ignore
        self._prs: Optional[Presentation] = None
        self._slide_count = 0

    def compose(self, pitch: Pitch, output_path: Optional[Path] = None) -> ComposerResult:
        """Generate a PowerPoint presentation from a pitch."""
        warnings: list[str] = []
        errors: list[str] = []

        try:
            # Create presentation
            self._prs = Presentation()
            self._slide_count = 0

            # Set slide dimensions
            self._prs.slide_width = Inches(self.config.slide_width)
            self._prs.slide_height = Inches(self.config.slide_height)

            # Get color palette
            palette = self.config.get_palette()

            # Create slides from presentation outline
            slides = pitch.to_presentation_outline()

            for slide_data in slides:
                slide_type = slide_data.get("slide_type", "content")

                if slide_type == "title":
                    self._create_title_slide(slide_data, palette, pitch)
                elif slide_type == "cta":
                    self._create_cta_slide(slide_data, palette)
                else:
                    self._create_content_slide(slide_data, palette, warnings)

            # Add elevator pitch slide at the end
            if pitch.elevator_pitch:
                self._create_elevator_pitch_slide(pitch, palette)

            # Add closing slide
            self._create_closing_slide(pitch, palette)

            # Resolve output path and save
            path = self._resolve_output_path(output_path, ".pptx")
            self._prs.save(str(path))

            # Get file size
            file_size = path.stat().st_size

            return ComposerResult(
                success=True,
                output_path=path,
                file_size_bytes=file_size,
                page_count=self._slide_count,
                warnings=warnings,
                metadata={
                    "format": "pptx",
                    "theme": self.config.theme.value,
                    "slide_dimensions": f"{self.config.slide_width}x{self.config.slide_height}",
                },
            )

        except Exception as e:
            errors.append(f"Failed to generate PPTX: {str(e)}")
            return ComposerResult(
                success=False,
                errors=errors,
                warnings=warnings,
            )

    def _create_title_slide(
        self, slide_data: dict, palette: dict[str, str], pitch: Pitch
    ) -> None:
        """Create the title slide."""
        slide_layout = self._prs.slide_layouts[6]  # Blank layout
        slide = self._prs.slides.add_slide(slide_layout)
        self._slide_count += 1

        # Add background
        self._add_gradient_background(slide, palette)

        # Title
        title_text = slide_data.get("title", pitch.title)
        # Truncate if too long
        if len(title_text) > 100:
            title_text = title_text[:97] + "..."

        title_box = slide.shapes.add_textbox(
            Inches(self.config.margin_left),
            Inches(2.5),
            Inches(self.config.slide_width - self.config.margin_left - self.config.margin_right),
            Inches(1.5),
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = title_text
        title_para.font.size = Pt(self.config.title_font_size)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(*self._hex_to_rgb(palette["text"]))
        title_para.font.name = self.config.title_font_name
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_text = slide_data.get("subtitle", pitch.subtitle or pitch.executive_summary)
        if subtitle_text:
            if len(subtitle_text) > 200:
                subtitle_text = subtitle_text[:197] + "..."

            subtitle_box = slide.shapes.add_textbox(
                Inches(self.config.margin_left),
                Inches(4.2),
                Inches(self.config.slide_width - self.config.margin_left - self.config.margin_right),
                Inches(1),
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = subtitle_text
            subtitle_para.font.size = Pt(self.config.subtitle_font_size)
            subtitle_para.font.color.rgb = RGBColor(*self._hex_to_rgb(palette["text_light"]))
            subtitle_para.font.name = self.config.body_font_name
            subtitle_para.alignment = PP_ALIGN.CENTER

        # Product name at bottom
        product_box = slide.shapes.add_textbox(
            Inches(self.config.margin_left),
            Inches(6.5),
            Inches(self.config.slide_width - self.config.margin_left - self.config.margin_right),
            Inches(0.5),
        )
        product_frame = product_box.text_frame
        product_para = product_frame.paragraphs[0]
        product_para.text = pitch.product_name
        product_para.font.size = Pt(14)
        product_para.font.color.rgb = RGBColor(*self._hex_to_rgb(palette["accent"]))
        product_para.font.name = self.config.body_font_name
        product_para.alignment = PP_ALIGN.CENTER

    def _create_content_slide(
        self, slide_data: dict, palette: dict[str, str], warnings: list[str]
    ) -> None:
        """Create a content slide with title and bullets."""
        slide_layout = self._prs.slide_layouts[6]  # Blank layout
        slide = self._prs.slides.add_slide(slide_layout)
        self._slide_count += 1

        # Add subtle background
        self._add_header_accent(slide, palette)

        # Title
        title_text = slide_data.get("title", "")
        title_box = slide.shapes.add_textbox(
            Inches(self.config.margin_left),
            Inches(self.config.margin_top),
            Inches(self.config.slide_width - self.config.margin_left - self.config.margin_right),
            Inches(0.8),
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = title_text
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(*self._hex_to_rgb(palette["primary"]))
        title_para.font.name = self.config.title_font_name

        # Content area
        bullets = slide_data.get("bullets", [])
        visual_assets = slide_data.get("visual_assets", [])

        # Determine layout based on visual assets
        has_image = any(
            va.get("type") == "image" and va.get("local_path")
            for va in visual_assets
        ) if self.config.include_visual_assets else False

        if has_image:
            # Two-column layout
            content_width = 6.5
            content_left = self.config.margin_left
            image_left = 7.5
            image_width = 5.0
        else:
            # Full-width layout
            content_width = self.config.slide_width - self.config.margin_left - self.config.margin_right
            content_left = self.config.margin_left

        # Add bullets
        content_top = 1.8
        content_height = 5.0

        if bullets:
            # Limit bullets
            display_bullets = bullets[: self.config.max_bullets_per_slide]
            if len(bullets) > self.config.max_bullets_per_slide:
                warnings.append(
                    f"Slide '{title_text}' has {len(bullets)} bullets, "
                    f"truncated to {self.config.max_bullets_per_slide}"
                )

            content_box = slide.shapes.add_textbox(
                Inches(content_left),
                Inches(content_top),
                Inches(content_width),
                Inches(content_height),
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for i, bullet in enumerate(display_bullets):
                if i == 0:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()

                # Handle long bullets - extract main point
                if len(bullet) > 200:
                    # Try to split at colon if present
                    if ":" in bullet[:80]:
                        bullet = bullet.split(":")[0] + ":"
                    else:
                        bullet = bullet[:197] + "..."

                para.text = f"• {bullet}"
                para.font.size = Pt(self.config.bullet_font_size)
                para.font.color.rgb = RGBColor(*self._hex_to_rgb(palette["text"]))
                para.font.name = self.config.body_font_name
                para.space_after = Pt(12)
                para.level = 0

        # Add image if available
        if has_image and self.config.include_visual_assets:
            for va in visual_assets:
                if va.get("type") == "image" and va.get("local_path"):
                    img_path = self._load_image_if_exists(va.get("local_path"))
                    if img_path:
                        try:
                            slide.shapes.add_picture(
                                str(img_path),
                                Inches(image_left),
                                Inches(content_top),
                                width=Inches(image_width),
                            )
                        except Exception as e:
                            warnings.append(f"Could not add image: {e}")
                    break  # Only add first image

        # Add speaker notes if configured
        if self.config.include_speaker_notes:
            notes = slide_data.get("notes", [])
            if notes:
                notes_slide = slide.notes_slide
                notes_frame = notes_slide.notes_text_frame
                notes_frame.text = "\n".join(f"• {note}" for note in notes)

        # Add slide number
        if self.config.add_slide_numbers:
            self._add_slide_number(slide, palette)

    def _create_cta_slide(self, slide_data: dict, palette: dict[str, str]) -> None:
        """Create a call-to-action slide."""
        slide_layout = self._prs.slide_layouts[6]  # Blank layout
        slide = self._prs.slides.add_slide(slide_layout)
        self._slide_count += 1

        # Add accent background
        self._add_cta_background(slide, palette)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(self.config.margin_left),
            Inches(1.5),
            Inches(self.config.slide_width - self.config.margin_left - self.config.margin_right),
            Inches(1),
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.get("title", "Next Steps")
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(*self._hex_to_rgb(palette["primary"]))
        title_para.font.name = self.config.title_font_name
        title_para.alignment = PP_ALIGN.CENTER

        # Primary CTA
        primary_cta = slide_data.get("primary_cta", "")
        if primary_cta:
            cta_box = slide.shapes.add_textbox(
                Inches(1.5),
                Inches(3),
                Inches(self.config.slide_width - 3),
                Inches(1.2),
            )
            cta_frame = cta_box.text_frame
            cta_frame.word_wrap = True
            cta_para = cta_frame.paragraphs[0]
            cta_para.text = primary_cta
            cta_para.font.size = Pt(24)
            cta_para.font.bold = True
            cta_para.font.color.rgb = RGBColor(*self._hex_to_rgb(palette["accent"]))
            cta_para.font.name = self.config.body_font_name
            cta_para.alignment = PP_ALIGN.CENTER

        # Next steps
        next_steps = slide_data.get("next_steps", [])
        if next_steps:
            steps_top = 4.5
            steps_box = slide.shapes.add_textbox(
                Inches(2),
                Inches(steps_top),
                Inches(self.config.slide_width - 4),
                Inches(2.5),
            )
            steps_frame = steps_box.text_frame
            steps_frame.word_wrap = True

            for i, step in enumerate(next_steps[:4]):  # Limit to 4 steps
                if i == 0:
                    para = steps_frame.paragraphs[0]
                else:
                    para = steps_frame.add_paragraph()

                para.text = f"{i + 1}. {step}"
                para.font.size = Pt(18)
                para.font.color.rgb = RGBColor(*self._hex_to_rgb(palette["text"]))
                para.font.name = self.config.body_font_name
                para.space_after = Pt(8)
                para.alignment = PP_ALIGN.LEFT

        if self.config.add_slide_numbers:
            self._add_slide_number(slide, palette)

    def _create_elevator_pitch_slide(self, pitch: Pitch, palette: dict[str, str]) -> None:
        """Create an elevator pitch summary slide."""
        slide_layout = self._prs.slide_layouts[6]
        slide = self._prs.slides.add_slide(slide_layout)
        self._slide_count += 1

        self._add_header_accent(slide, palette)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(self.config.margin_left),
            Inches(self.config.margin_top),
            Inches(self.config.slide_width - self.config.margin_left - self.config.margin_right),
            Inches(0.8),
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "The 30-Second Pitch"
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(*self._hex_to_rgb(palette["primary"]))
        title_para.font.name = self.config.title_font_name

        # Elevator pitch content
        content_box = slide.shapes.add_textbox(
            Inches(1),
            Inches(2),
            Inches(self.config.slide_width - 2),
            Inches(4),
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_para = content_frame.paragraphs[0]
        content_para.text = f'"{pitch.elevator_pitch}"'
        content_para.font.size = Pt(22)
        content_para.font.italic = True
        content_para.font.color.rgb = RGBColor(*self._hex_to_rgb(palette["text"]))
        content_para.font.name = self.config.body_font_name
        content_para.alignment = PP_ALIGN.CENTER

        # Key messages
        if pitch.key_messages:
            messages_top = 5.0
            messages_box = slide.shapes.add_textbox(
                Inches(1.5),
                Inches(messages_top),
                Inches(self.config.slide_width - 3),
                Inches(2),
            )
            messages_frame = messages_box.text_frame
            messages_frame.word_wrap = True

            header_para = messages_frame.paragraphs[0]
            header_para.text = "Key Messages:"
            header_para.font.size = Pt(14)
            header_para.font.bold = True
            header_para.font.color.rgb = RGBColor(*self._hex_to_rgb(palette["text_light"]))
            header_para.font.name = self.config.body_font_name

            for msg in pitch.key_messages[:5]:
                para = messages_frame.add_paragraph()
                para.text = f"• {msg}"
                para.font.size = Pt(12)
                para.font.color.rgb = RGBColor(*self._hex_to_rgb(palette["text_light"]))
                para.font.name = self.config.body_font_name

        if self.config.add_slide_numbers:
            self._add_slide_number(slide, palette)

    def _create_closing_slide(self, pitch: Pitch, palette: dict[str, str]) -> None:
        """Create a closing/thank you slide."""
        slide_layout = self._prs.slide_layouts[6]
        slide = self._prs.slides.add_slide(slide_layout)
        self._slide_count += 1

        self._add_gradient_background(slide, palette)

        # Thank you text
        thanks_box = slide.shapes.add_textbox(
            Inches(self.config.margin_left),
            Inches(2.5),
            Inches(self.config.slide_width - self.config.margin_left - self.config.margin_right),
            Inches(1.5),
        )
        thanks_frame = thanks_box.text_frame
        thanks_para = thanks_frame.paragraphs[0]
        thanks_para.text = "Thank You"
        thanks_para.font.size = Pt(48)
        thanks_para.font.bold = True
        thanks_para.font.color.rgb = RGBColor(*self._hex_to_rgb(palette["text"]))
        thanks_para.font.name = self.config.title_font_name
        thanks_para.alignment = PP_ALIGN.CENTER

        # One-liner
        if pitch.one_liner:
            liner_box = slide.shapes.add_textbox(
                Inches(1),
                Inches(4),
                Inches(self.config.slide_width - 2),
                Inches(1),
            )
            liner_frame = liner_box.text_frame
            liner_frame.word_wrap = True
            liner_para = liner_frame.paragraphs[0]
            liner_para.text = pitch.one_liner
            liner_para.font.size = Pt(20)
            liner_para.font.color.rgb = RGBColor(*self._hex_to_rgb(palette["text_light"]))
            liner_para.font.name = self.config.body_font_name
            liner_para.alignment = PP_ALIGN.CENTER

        # Product URL
        url_box = slide.shapes.add_textbox(
            Inches(self.config.margin_left),
            Inches(6),
            Inches(self.config.slide_width - self.config.margin_left - self.config.margin_right),
            Inches(0.5),
        )
        url_frame = url_box.text_frame
        url_para = url_frame.paragraphs[0]
        url_para.text = pitch.product_url
        url_para.font.size = Pt(14)
        url_para.font.color.rgb = RGBColor(*self._hex_to_rgb(palette["accent"]))
        url_para.font.name = self.config.body_font_name
        url_para.alignment = PP_ALIGN.CENTER

    def _add_gradient_background(self, slide, palette: dict[str, str]) -> None:
        """Add a gradient-like background using a rectangle."""
        # Add a subtle colored rectangle at the top
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(self.config.slide_width),
            Inches(0.15),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(palette["primary"]))
        shape.line.fill.background()

    def _add_header_accent(self, slide, palette: dict[str, str]) -> None:
        """Add a subtle header accent line."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(self.config.margin_left),
            Inches(1.5),
            Inches(2),
            Inches(0.05),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(palette["accent"]))
        shape.line.fill.background()

    def _add_cta_background(self, slide, palette: dict[str, str]) -> None:
        """Add background styling for CTA slide."""
        # Top accent bar
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(self.config.slide_width),
            Inches(0.25),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(palette["accent"]))
        shape.line.fill.background()

    def _add_slide_number(self, slide, palette: dict[str, str]) -> None:
        """Add slide number to bottom right."""
        num_box = slide.shapes.add_textbox(
            Inches(self.config.slide_width - 1),
            Inches(self.config.slide_height - 0.4),
            Inches(0.5),
            Inches(0.3),
        )
        num_frame = num_box.text_frame
        num_para = num_frame.paragraphs[0]
        num_para.text = str(self._slide_count)
        num_para.font.size = Pt(10)
        num_para.font.color.rgb = RGBColor(*self._hex_to_rgb(palette["text_light"]))
        num_para.font.name = self.config.body_font_name
        num_para.alignment = PP_ALIGN.RIGHT
